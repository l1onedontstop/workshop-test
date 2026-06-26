"""Generate SparkForge product introduction PPT"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Color Palette ──
BG_DARK = RGBColor(0x0F, 0x0F, 0x13)
BG_CARD = RGBColor(0x1A, 0x1A, 0x2E)
BRAND = RGBColor(0x7C, 0x5C, 0xFC)
ACCENT = RGBColor(0x00, 0xD4, 0xFF)
GREEN = RGBColor(0x00, 0xE6, 0x76)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0xA0, 0xA0, 0xB0)
YELLOW = RGBColor(0xFF, 0xD7, 0x40)
RED = RGBColor(0xFF, 0x52, 0x52)
ORANGE = RGBColor(0xFF, 0x85, 0x00)
PURPLE = RGBColor(0xA0, 0x70, 0xFF)
DARK_BORDER = RGBColor(0x2A, 0x2A, 0x40)
FONT = 'Microsoft YaHei'

def dark_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG_DARK
    return slide

def card(slide, left, top, width, height, color=BG_CARD, border=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if border:
        shape.line.color.rgb = border
        shape.line.width = Pt(1)
    return shape

def txt(slide, left, top, width, height, text, size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = FONT
    p.alignment = align
    return txBox

def multiline(slide, left, top, width, height, lines, size=16, color=WHITE, spacing=1.5):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = FONT
        p.space_after = Pt(size * (spacing - 1))
    return txBox

def foot(slide, page):
    txt(slide, Inches(0.5), Inches(6.9), Inches(12), Inches(0.4),
        f'SparkForge v1.0.0  |  个人影响力锻造工坊  |  {page}',
        size=10, color=GRAY, align=PP_ALIGN.CENTER)

# ═══════════════════ Slide 1: Cover ═══════════════════
s = dark_slide()
bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.06))
bar.fill.solid(); bar.fill.fore_color.rgb = BRAND; bar.line.fill.background()
txt(s, Inches(1), Inches(1.3), Inches(11), Inches(0.8), 'SparkForge', size=60, color=WHITE, bold=True)
txt(s, Inches(1), Inches(2.2), Inches(11), Inches(0.5), '个人影响力锻造工坊', size=36, color=ACCENT)
txt(s, Inches(1), Inches(3.0), Inches(11), Inches(0.4), 'AI 驱动的短视频内容创作全流程平台', size=22, color=GRAY)

pills = ['AI脚本生成', '7维智能评分', '热点匹配', '数据复盘', 'IP蓝图策略']
for i, pill in enumerate(pills):
    x = Inches(1 + i * 2.3)
    card(s, x, Inches(4.0), Inches(2.1), Inches(0.6), BG_CARD, DARK_BORDER)
    txt(s, x + Inches(0.15), Inches(4.15), Inches(1.8), Inches(0.4), pill, size=14, color=WHITE, align=PP_ALIGN.CENTER)

txt(s, Inches(1), Inches(5.3), Inches(11), Inches(0.4), 'Electron + React 19 + TypeScript + Tailwind CSS + AI', size=14, color=GRAY, align=PP_ALIGN.CENTER)
foot(s, '01 / 10')

# ═══════════════════ Slide 2: Problem & Solution ═══════════════════
s = dark_slide()
txt(s, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6), '创作者面临的困境 vs SparkForge 解决方案', size=32, color=WHITE, bold=True)

problems = [
    ('不会写口播文案', '对着镜头不知道说什么，写出来的像念稿'),
    ('不知道好不好', '写完脚本心里没底，发了才知道效果'),
    ('优化无方向', '改了又改，不知道哪里需要改、怎么改'),
    ('数据不反馈', '发了就发了，播放数据不会反哺到下一次创作'),
    ('选题靠运气', '不知道最近什么话题火，凭感觉选方向'),
]
for i, (title, desc) in enumerate(problems):
    y = Inches(1.3 + i * 1.05)
    card(s, Inches(0.5), y, Inches(5.8), Inches(0.85), BG_CARD, DARK_BORDER)
    txt(s, Inches(0.7), y + Inches(0.08), Inches(5.4), Inches(0.3), f'❌ {title}', size=16, color=RED, bold=True)
    txt(s, Inches(0.7), y + Inches(0.42), Inches(5.4), Inches(0.3), desc, size=12, color=GRAY)

solutions = [
    ('AI 脚本生成', '输入话题，AI 自动生成完整口播文案 + 8段拍摄方案'),
    ('7 维智能评分', '开篇钩子·叙事节奏·观点锐度·实用密度·情绪共鸣·结构完整·表达效果'),
    ('AI 优化闭环', '评分定位弱项 → AI定向优化 → 重新评分 → 得分轨迹追踪'),
    ('数据复盘进化', '发布数据回传 → 预测偏差分析 → 评分规则自动进化'),
    ('热点实时匹配', '全网热榜监控 + AI 匹配你的行业 → 选题不再靠猜'),
]
for i, (title, desc) in enumerate(solutions):
    y = Inches(1.3 + i * 1.05)
    card(s, Inches(6.8), y, Inches(5.8), Inches(0.85), BG_CARD, DARK_BORDER)
    txt(s, Inches(7.0), y + Inches(0.08), Inches(5.4), Inches(0.3), f'✅ {title}', size=16, color=GREEN, bold=True)
    txt(s, Inches(7.0), y + Inches(0.42), Inches(5.4), Inches(0.3), desc, size=12, color=GRAY)

foot(s, '02 / 10')

# ═══════════════════ Slide 3: Core Workflow ═══════════════════
s = dark_slide()
txt(s, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6), '核心创作闭环', size=36, color=WHITE, bold=True)

steps = [
    ('IP蓝图', '定义人设', BRAND),
    ('选题策划', '热点匹配', ACCENT),
    ('AI写稿', '口播+方案', GREEN),
    ('智能评分', '7维打分', YELLOW),
    ('AI优化', '定向改写', RED),
    ('发布追踪', '多平台同步', ACCENT),
    ('数据复盘', '规则进化', BRAND),
]
for i, (title, desc, color) in enumerate(steps):
    x = Inches(0.3 + i * 1.82)
    card(s, x, Inches(1.3), Inches(1.65), Inches(2.2), BG_CARD, color)
    txt(s, x + Inches(0.1), Inches(1.5), Inches(1.45), Inches(0.4), str(i+1), size=32, color=color, bold=True, align=PP_ALIGN.CENTER)
    txt(s, x + Inches(0.1), Inches(2.0), Inches(1.45), Inches(0.35), title, size=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    txt(s, x + Inches(0.1), Inches(2.4), Inches(1.45), Inches(0.3), desc, size=11, color=GRAY, align=PP_ALIGN.CENTER)
    if i < len(steps) - 1:
        txt(s, x + Inches(1.6), Inches(2.1), Inches(0.25), Inches(0.3), '>', size=16, color=GRAY)

txt(s, Inches(0.8), Inches(3.9), Inches(11), Inches(0.5), '核心指标', size=24, color=WHITE, bold=True)
metrics = [
    ('7', '评分维度'),
    ('8', '方案段落'),
    ('25+', '校准样本'),
    ('15', '功能页面'),
    ('2', '平台支持'),
]
for i, (num, label) in enumerate(metrics):
    x = Inches(0.8 + i * 2.5)
    card(s, x, Inches(4.6), Inches(2.2), Inches(1.4), BG_CARD, DARK_BORDER)
    txt(s, x + Inches(0.1), Inches(4.75), Inches(2.0), Inches(0.5), num, size=30, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
    txt(s, x + Inches(0.1), Inches(5.35), Inches(2.0), Inches(0.35), label, size=14, color=GRAY, align=PP_ALIGN.CENTER)

# Architecture note
txt(s, Inches(0.8), Inches(6.3), Inches(11), Inches(0.4), 'Electron桌面应用 + React前端 + TypeScript全栈 + 多AI引擎支持', size=13, color=GRAY)
foot(s, '03 / 10')

# ═══════════════════ Slide 4: AI Script ═══════════════════
s = dark_slide()
txt(s, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6), 'AI 脚本生成 —— 8段完整拍摄方案', size=36, color=WHITE, bold=True)

sections = [
    ('口播文案', '完整的口播脚本，口语化表达，像在说话而非念稿'),
    ('风格定义', '视频风格、调色方向、情绪基调的统一规划'),
    ('分镜脚本', '镜头拆解、画面描述、时长分配的详细方案'),
    ('设备建议', '拍摄设备、灯光、收音方案的具体推荐'),
    ('场景布置', '拍摄场景、道具、背景的完整建议'),
    ('后期制作', '剪辑节奏、特效风格、字幕方案的指导'),
    ('封面设计', '封面图方案、标题文案的优化建议'),
    ('AI自评分', '预估各维度得分 + 总分，提前知道脚本质量'),
]
for i, (title, desc) in enumerate(sections):
    col = i % 2
    row = i // 2
    x = Inches(0.5 + col * 6.3)
    y = Inches(1.3 + row * 1.4)
    card(s, x, y, Inches(5.9), Inches(1.15), BG_CARD, DARK_BORDER)
    txt(s, x + Inches(0.2), y + Inches(0.08), Inches(1.5), Inches(0.3), f'0{i+1}', size=28, color=BRAND, bold=True)
    txt(s, x + Inches(1.0), y + Inches(0.08), Inches(4.7), Inches(0.3), title, size=19, color=WHITE, bold=True)
    txt(s, x + Inches(1.0), y + Inches(0.5), Inches(4.7), Inches(0.4), desc, size=12, color=GRAY)

txt(s, Inches(0.8), Inches(7.0), Inches(11), Inches(0.3), '一键生成 → 自动保存 → 直接进入评分优化流程', size=14, color=ACCENT)
foot(s, '04 / 10')

# ═══════════════════ Slide 5: 7-D Scoring ═══════════════════
s = dark_slide()
txt(s, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6), '7 维智能评分系统', size=36, color=WHITE, bold=True)
txt(s, Inches(0.8), Inches(0.9), Inches(11), Inches(0.3), '25+ 样本拟合 · 随复盘数据自动进化 · AI内容教练', size=15, color=GRAY)

dims = [
    ('开篇钩子', '20%', '前3秒能不能让人停下来', BRAND),
    ('叙事节奏', '15%', '信息密度和情绪起伏', ACCENT),
    ('观点锐度', '15%', '有没有让人卧槽的洞见', GREEN),
    ('实用密度', '15%', '观众能拿走什么', YELLOW),
    ('情绪共鸣', '15%', '会不会想转发/评论', RED),
    ('结构完整', '10%', '开头-展开-高潮-结尾', ORANGE),
    ('表达效果', '10%', '口语化、画面感、适合口播', PURPLE),
]
for i, (name, weight, desc, color) in enumerate(dims):
    x = Inches(0.3 + i * 1.82)
    card(s, x, Inches(1.5), Inches(1.65), Inches(3.0), BG_CARD, color)
    txt(s, x + Inches(0.1), Inches(1.7), Inches(1.45), Inches(0.3), f'{i+1}', size=24, color=color, bold=True, align=PP_ALIGN.CENTER)
    txt(s, x + Inches(0.1), Inches(2.1), Inches(1.45), Inches(0.3), name, size=17, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    txt(s, x + Inches(0.1), Inches(2.5), Inches(1.45), Inches(0.3), weight, size=22, color=color, bold=True, align=PP_ALIGN.CENTER)
    txt(s, x + Inches(0.1), Inches(2.9), Inches(1.45), Inches(0.5), desc, size=11, color=GRAY, align=PP_ALIGN.CENTER)

txt(s, Inches(0.8), Inches(4.8), Inches(11), Inches(0.5), '校准闭环', size=24, color=WHITE, bold=True)
card(s, Inches(0.8), Inches(5.4), Inches(11.7), Inches(1.1), BG_CARD, DARK_BORDER)
calib = '生成脚本 → AI评分 → 发布追踪 → 实际数据回传 → 预测 vs 实际偏差分析 → 评分权重自动调整 → 越来越准'
txt(s, Inches(1.2), Inches(5.6), Inches(11), Inches(0.5), calib, size=16, color=ACCENT, align=PP_ALIGN.CENTER)
txt(s, Inches(1.2), Inches(6.1), Inches(11), Inches(0.3), '每次复盘积累数据，让评分规则随你的创作风格自动进化', size=13, color=GRAY, align=PP_ALIGN.CENTER)
foot(s, '05 / 10')

# ═══════════════════ Slide 6: AI Optimize ═══════════════════
s = dark_slide()
txt(s, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6), 'AI 优化闭环', size=36, color=WHITE, bold=True)

flow = [
    ('生成脚本', 'AI自动生成口播文案', BRAND),
    ('智能评分', '7维打分，定位弱项', ACCENT),
    ('AI优化', '根据弱项定向改写', GREEN),
    ('对比审核', '原文 vs 优化版左右对比', YELLOW),
    ('接受打分', '自动重新评分', RED),
    ('追踪进化', '得分轨迹 6.5 -> 7.2 -> 7.8', BRAND),
]
for i, (title, desc, color) in enumerate(flow):
    x = Inches(0.3 + i * 2.15)
    card(s, x, Inches(1.3), Inches(1.95), Inches(2.0), BG_CARD, color)
    txt(s, x + Inches(0.1), Inches(1.5), Inches(1.75), Inches(0.4), str(i+1), size=32, color=color, bold=True, align=PP_ALIGN.CENTER)
    txt(s, x + Inches(0.1), Inches(2.0), Inches(1.75), Inches(0.35), title, size=15, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    txt(s, x + Inches(0.1), Inches(2.4), Inches(1.75), Inches(0.4), desc, size=11, color=GRAY, align=PP_ALIGN.CENTER)
    if i < len(flow) - 1:
        txt(s, x + Inches(1.9), Inches(2.0), Inches(0.25), Inches(0.3), '>', size=16, color=GRAY)

txt(s, Inches(0.8), Inches(3.7), Inches(11), Inches(0.5), '优化功能亮点', size=24, color=WHITE, bold=True)
highlights = [
    '得分轨迹：记录每次优化后的评分变化, 可视化进步过程 (例: 6.5 -> 7.2 -> 7.8)',
    '一键撤销：不满意可以回退到上一版, 保留完整优化历史, 随时对比各版本',
    '过期标记：修改脚本后评分自动标记"已修改", 避免陈旧评分误导判断',
    '自动保存：0.5秒防抖自动保存, 不再担心内容丢失',
    '改动摘要：显示每次优化的字数变化和段落调整, 量化优化效果',
    '最低分高亮：自动将最低维度排在最前, <=5分红色高亮 + 改善提示',
]
multiline(s, Inches(1.2), Inches(4.3), Inches(11), Inches(2.2), highlights, size=15, color=GRAY, spacing=1.7)
foot(s, '06 / 10')

# ═══════════════════ Slide 7: Features ═══════════════════
s = dark_slide()
txt(s, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6), '功能全景', size=36, color=WHITE, bold=True)

features = [
    ('工作台', '项目管理 · 快捷操作 · AI Coach提示', BRAND),
    ('IP蓝图', '6问引导 · 策略定义 · 对标分析 · 自动转化', ACCENT),
    ('选题池', '选题管理 · 标签分类 · 优先级排序 · 状态跟踪', GREEN),
    ('热点匹配', '实时热榜 · AI行业匹配 · 灵感推荐 · 一键转化', YELLOW),
    ('脚本编辑器', 'AI生成 · 7维评分 · 优化闭环 · 一键撤销 · 8段方案', RED),
    ('数据复盘', '预测偏差分析 · 规则自动进化 · 历史对比 · 学习积累', ORANGE),
    ('方案管理', '拍摄方案列表 · 批量生成 · 脚本联动 · 状态管理', PURPLE),
    ('发布管理', '多平台追踪 · 数据回传 · 效果分析 · 发布时间记录', ACCENT),
    ('受众画像', '目标受众定义 · 需求分析 · 内容偏好匹配', BRAND),
    ('对标分析', '竞品研究 · 差异定位 · 策略建议 · 数据对比', GREEN),
]
for i, (name, desc, color) in enumerate(features):
    col = i % 2
    row = i // 2
    x = Inches(0.5 + col * 6.3)
    y = Inches(1.2 + row * 1.15)
    card(s, x, y, Inches(5.9), Inches(0.95), BG_CARD, color)
    txt(s, x + Inches(0.2), y + Inches(0.08), Inches(5.5), Inches(0.3), name, size=18, color=WHITE, bold=True)
    txt(s, x + Inches(0.2), y + Inches(0.45), Inches(5.5), Inches(0.3), desc, size=12, color=GRAY)

foot(s, '07 / 10')

# ═══════════════════ Slide 8: Tech ═══════════════════
s = dark_slide()
txt(s, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6), '技术架构', size=36, color=WHITE, bold=True)

layers = [
    ('桌面客户端', ['Electron 33 框架', 'macOS + Windows 双平台', '原生TitleBar + 系统集成', 'Electron-Builder 打包分发'], BRAND),
    ('前端渲染层', ['React 19 + TypeScript 5.7', 'Tailwind CSS 3.4 暗色主题', 'Zustand 5 状态管理', 'Radix UI + Lucide 图标', 'CVA 组件变体管理'], ACCENT),
    ('IPC 通信层', ['contextBridge 安全桥接', 'preload 类型安全API', '文件系统操作', '进程管理 + 系统调用'], GREEN),
    ('AI 服务层', ['多引擎支持 (可切换)', 'OpenAI / DeepSeek / 等', '请求队列 + 速率限制', '智能重试 + 错误翻译'], YELLOW),
    ('数据存储层', ['Markdown + JSON 文件系统', 'SQLite 向量数据库 (HNSW)', '自动备份 + 版本管理', '离线优先 + 本地优先'], RED),
]
for i, (name, items, color) in enumerate(layers):
    y = Inches(1.2 + i * 1.15)
    card(s, Inches(0.5), y, Inches(3.0), Inches(0.95), color, None)
    txt(s, Inches(0.7), y + Inches(0.3), Inches(2.6), Inches(0.35), name, size=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    for j, item in enumerate(items):
        tx = Inches(3.8 + (j % 3) * 3.1)
        ty = y + Inches(0.1 if j < 3 else 0.5)
        txt(s, tx, ty, Inches(2.9), Inches(0.25), f'  {item}', size=11, color=GRAY)

txt(s, Inches(0.8), Inches(7.0), Inches(11), Inches(0.3), 'TypeScript 全栈 · 共享类型系统 · useReducer 状态管理 · @common 共享常量', size=13, color=GRAY, align=PP_ALIGN.CENTER)
foot(s, '08 / 10')

# ═══════════════════ Slide 9: Data Flow ═══════════════════
s = dark_slide()
txt(s, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6), '数据闭环 & 进化机制', size=36, color=WHITE, bold=True)

txt(s, Inches(0.8), Inches(1.2), Inches(5), Inches(0.4), '内容质量进化', size=22, color=ACCENT, bold=True)
evo_steps = [
    '1. AI 生成脚本 + 预测评分',
    '2. 创作者修改 + AI 优化 + 重新评分',
    '3. 发布到各平台，收集真实数据',
    '4. 实际数据 vs 预测评分 偏差分析',
    '5. 评分权重根据偏差自动调整',
    '6. 积累越多数据，评分越准确',
]
multiline(s, Inches(0.8), Inches(1.8), Inches(5.5), Inches(2.8), evo_steps, size=15, color=GRAY, spacing=1.7)

txt(s, Inches(7.2), Inches(1.2), Inches(5), Inches(0.4), '技术特性', size=22, color=GREEN, bold=True)
tech_highlights = [
    '本地优先：所有数据存储在本地，不上传云端',
    '离线可用：无需联网即可编辑和管理内容',
    '多AI引擎：支持切换不同AI服务商',
    '自动备份：定时自动备份项目数据',
    '类型安全：TypeScript全栈，编译时发现错误',
    '组件化：共享UI组件库，一致的产品体验',
]
multiline(s, Inches(7.2), Inches(1.8), Inches(5.5), Inches(2.8), tech_highlights, size=15, color=GRAY, spacing=1.7)

# Use cases
txt(s, Inches(0.8), Inches(4.8), Inches(11), Inches(0.4), '适用场景', size=22, color=YELLOW, bold=True)
scenarios = [
    ('短视频创作者', '系统化管理内容创作，AI辅助提升脚本质量'),
    ('知识博主', 'IP定位策略 + 选题规划 + 专业评分'),
    ('内容团队', '多人协作方案管理 + 数据驱动的质量进化'),
    ('个人IP孵化', '从0到1建立内容体系 + AI全程辅助'),
]
for i, (who, what) in enumerate(scenarios):
    x = Inches(0.5 + i * 3.15)
    card(s, x, Inches(5.4), Inches(2.9), Inches(1.2), BG_CARD, DARK_BORDER)
    txt(s, x + Inches(0.15), Inches(5.55), Inches(2.6), Inches(0.3), who, size=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    txt(s, x + Inches(0.15), Inches(5.9), Inches(2.6), Inches(0.5), what, size=12, color=GRAY, align=PP_ALIGN.CENTER)

foot(s, '09 / 10')

# ═══════════════════ Slide 10: CTA ═══════════════════
s = dark_slide()
bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.06))
bar.fill.solid(); bar.fill.fore_color.rgb = BRAND; bar.line.fill.background()

txt(s, Inches(1), Inches(1.0), Inches(11), Inches(0.8), 'SparkForge', size=56, color=WHITE, bold=True)
txt(s, Inches(1), Inches(1.8), Inches(11), Inches(0.5), '让每个人都能专业地做内容', size=28, color=ACCENT)

pillars = [
    ('AI 赋能', 'AI 不是替代创作者，而是放大每个人的创作力'),
    ('数据驱动', '评分 + 复盘 + 进化，让创作从玄学变成科学'),
    ('完整闭环', '从灵感到发布到复盘，一站式创作工作流'),
]
for i, (title, desc) in enumerate(pillars):
    x = Inches(0.8 + i * 4.1)
    card(s, x, Inches(2.8), Inches(3.7), Inches(1.8), BG_CARD, DARK_BORDER)
    txt(s, x + Inches(0.2), Inches(3.0), Inches(3.3), Inches(0.5), title, size=24, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    txt(s, x + Inches(0.2), Inches(3.6), Inches(3.3), Inches(0.7), desc, size=14, color=GRAY, align=PP_ALIGN.CENTER)

txt(s, Inches(1), Inches(5.1), Inches(11), Inches(0.3), '技术栈', size=16, color=GRAY)
techs = ['Electron 33', 'React 19', 'TypeScript 5.7', 'Tailwind CSS 3.4', 'Zustand', 'Radix UI']
for i, tech in enumerate(techs):
    x = Inches(0.8 + i * 2.05)
    card(s, x, Inches(5.6), Inches(1.85), Inches(0.5), BG_CARD, DARK_BORDER)
    txt(s, x + Inches(0.1), Inches(5.7), Inches(1.65), Inches(0.3), tech, size=13, color=ACCENT, align=PP_ALIGN.CENTER)

txt(s, Inches(1), Inches(6.4), Inches(11), Inches(0.3), 'v1.0.0  |  macOS + Windows  |  github.com/l1onedontstop/workshop-test', size=12, color=GRAY, align=PP_ALIGN.CENTER)
foot(s, '10 / 10')

# ── Save ──
output_path = 'F:/ip工坊/ip-studio/SparkForge_产品介绍.pptx'
prs.save(output_path)
print(f'PPT saved: {output_path}')
print(f'{len(prs.slides)} slides created')
